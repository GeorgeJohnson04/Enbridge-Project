"""
Build the interpretive PowerPoint deck for the four OLS oil-price models.

Output: output/Oil_OLS_Model_Presentation.pptx
"""
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

import config

# -- Palette ----------------------------------------------------------------
INK     = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
CREAM   = RGBColor(0xF5, 0xF1, 0xE8)
AMBER   = RGBColor(0xE8, 0x89, 0x3E)
TEAL    = RGBColor(0x16, 0x69, 0x7A)
RED     = RGBColor(0xA8, 0x32, 0x4E)
GRAY    = RGBColor(0x55, 0x55, 0x60)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTGR = RGBColor(0xE0, 0xDC, 0xD2)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

CHARTS = config.CHARTS_DIR
OUTFILE = config.OUTPUT_DIR / "Oil_OLS_Model_Presentation.pptx"

# 16:9
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


# -- helpers ----------------------------------------------------------------
def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_left_bar(slide, color=AMBER, w=Inches(0.18)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, w, SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=INK, font=BODY_FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                size=14, color=INK, font=BODY_FONT, line_spacing=1.2,
                bullet_color=AMBER, bullet_char="\u25CF"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet
        b = p.add_run()
        b.text = f"{bullet_char}  "
        b.font.name = font; b.font.size = Pt(size); b.font.bold = True
        b.font.color.rgb = bullet_color
        # text
        if isinstance(item, tuple) and len(item) == 2:
            head, tail = item
            r1 = p.add_run(); r1.text = head
            r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run(); r2.text = " — " + tail
            r2.font.name = font; r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run(); r.text = item
            r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color
    return tb


def add_image(slide, path: Path, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_slide_number(slide, n, total):
    add_text(slide, Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.3),
             f"{n} / {total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def page_header(slide, title, kicker=None):
    """Standard content-slide header: small kicker above big title."""
    add_left_bar(slide)
    if kicker:
        add_text(slide, Inches(0.55), Inches(0.32), Inches(8), Inches(0.3),
                 kicker.upper(), size=11, bold=True, color=AMBER, font=BODY_FONT)
    add_text(slide, Inches(0.55), Inches(0.55), Inches(12.5), Inches(0.7),
             title, size=28, bold=True, color=INK, font=HEAD_FONT)


# -- Slides -----------------------------------------------------------------
def slide_title(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(s, INK)

    # Amber strip on the left
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = AMBER; band.line.fill.background()

    # Kicker
    add_text(s, Inches(0.9), Inches(2.0), Inches(11), Inches(0.4),
             "ENBRIDGE  |  OIL PRICE FORECAST",
             size=14, bold=True, color=AMBER, font=BODY_FONT)
    # Title
    add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.0),
             "Where Is Oil\nHeaded?",
             size=54, bold=True, color=WHITE, font=HEAD_FONT)
    # Subtitle
    add_text(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.6),
             "A 4-month price outlook through August 2026, built on 23 years of data",
             size=20, italic=True, color=CREAM, font=BODY_FONT)
    # Footer
    add_text(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.3),
             "Monthly data 2003 \u2013 2026  |  Forecast horizon: August 2026",
             size=12, color=LIGHTGR, font=BODY_FONT)


def _read_forecast_summary():
    """Returns dict with last actuals and August-2026 forecast values."""
    fcst = pd.read_csv(config.OUTPUT_DIR / "forecast.csv",
                       index_col="date", parse_dates=True)
    df = pd.read_csv(config.FEATURES_CSV, index_col="date", parse_dates=True)
    light_last = float(df["wti_real"].dropna().iloc[-1])
    heavy_last = float(df["wcs_real"].dropna().iloc[-1])
    light_end = float(fcst["light_price_pred"].iloc[-1])
    heavy_end = float(fcst["heavy_price_pred"].iloc[-1])
    light_lo  = float(fcst["light_price_lo_95"].iloc[-1])
    light_hi  = float(fcst["light_price_hi_95"].iloc[-1])
    heavy_lo  = float(fcst["heavy_price_lo_95"].iloc[-1])
    heavy_hi  = float(fcst["heavy_price_hi_95"].iloc[-1])
    return dict(
        light_last=light_last, heavy_last=heavy_last,
        light_end=light_end, heavy_end=heavy_end,
        light_lo=light_lo, light_hi=light_hi,
        heavy_lo=heavy_lo, heavy_hi=heavy_hi,
    )


def slide_executive_summary(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, "The forecast in one slide", kicker="The big picture")

    f = _read_forecast_summary()

    # Two big forecast tiles: Light and Heavy projection
    tiles = [
        ("Light oil (WTI)", f"${f['light_end']:.0f}",
            f"from ${f['light_last']:.0f} today  |  range ${f['light_lo']:.0f} \u2013 ${f['light_hi']:.0f}",
            TEAL),
        ("Heavy oil",       f"${f['heavy_end']:.0f}",
            f"from ${f['heavy_last']:.0f} today  |  range ${f['heavy_lo']:.0f} \u2013 ${f['heavy_hi']:.0f}",
            AMBER),
    ]
    tile_w, tile_h = Inches(6.05), Inches(2.0)
    gap = Inches(0.3)
    start_x = Inches(0.55)
    y0 = Inches(1.5)
    for i, (name, val, sub, accent) in enumerate(tiles):
        x = start_x + (tile_w + gap) * i
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, tile_w, tile_h)
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = accent; box.line.width = Pt(2)
        box.shadow.inherit = False
        add_text(s, x, y0 + Inches(0.2), tile_w, Inches(0.4),
                 name, size=14, bold=True, color=accent, font=BODY_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x, y0 + Inches(0.55), tile_w, Inches(1.1),
                 val, size=64, bold=True, color=INK, font=HEAD_FONT, align=PP_ALIGN.CENTER)
        add_text(s, x, y0 + Inches(1.6), tile_w, Inches(0.35),
                 "by August 2026  |  per barrel (2025 dollars)",
                 size=10, color=GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Sub-caption under tiles
    add_text(s, Inches(0.55), Inches(3.6), Inches(12.4), Inches(0.4),
             f"Today's prices: Light ${f['light_last']:.0f}  |  Heavy ${f['heavy_last']:.0f}.  "
             f"Both prices forecast to settle below current levels by next fiscal year-end.",
             size=12, italic=True, color=GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER)

    # Body interpretation
    add_text(s, Inches(0.55), Inches(4.15), Inches(12.4), Inches(0.45),
             "What this deck shows", size=18, bold=True, color=INK, font=HEAD_FONT)
    add_bullets(s, Inches(0.55), Inches(4.6), Inches(12.4), Inches(2.5), [
        ("Forecast through August 2026",
            "4-month projection for both light and heavy oil, with confidence ranges"),
        ("Built on 23 years of monthly data",
            "supply, demand, dollar strength, geopolitical risk, and conflict history since 2003"),
        ("Four supporting models behind the forecast",
            "split by oil type (light vs heavy) and conflict length (short vs long); each explains 92\u201396% of past prices"),
        ("Everything is reproducible from raw data",
            "one command refreshes the data, refits the models, and regenerates the forecast"),
    ], size=14, color=INK)
    add_slide_number(s, n, total)


def slide_methodology(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, "How the model works", kicker="The approach")

    # Left: the 2x2 segmentation grid
    left_x = Inches(0.55)
    top_y = Inches(1.5)
    cell_w, cell_h = Inches(2.55), Inches(1.45)
    add_text(s, left_x, top_y, Inches(5.5), Inches(0.4),
             "We split history into four cases",
             size=15, bold=True, color=INK, font=HEAD_FONT)

    # corner labels
    add_text(s, left_x + Inches(0.6), top_y + Inches(0.5), cell_w, Inches(0.3),
             "Light oil (WTI)", size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s, left_x + Inches(0.6) + cell_w + Inches(0.1), top_y + Inches(0.5), cell_w, Inches(0.3),
             "Heavy oil (WCS)", size=11, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # row labels and tiles
    row_labels = ["Short\nconflict\n(< 6 mo)", "Long\nconflict\n(6+ mo)"]
    cell_specs = [
        ("Model 1", "Light  |  Short", TEAL),
        ("Model 3", "Heavy  |  Short", AMBER),
        ("Model 2", "Light  |  Long", TEAL),
        ("Model 4", "Heavy  |  Long", AMBER),
    ]
    for r in range(2):
        ry = top_y + Inches(0.85) + r * (cell_h + Inches(0.12))
        add_text(s, left_x, ry, Inches(0.6), cell_h,
                 row_labels[r], size=10, bold=True, color=GRAY,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        for c in range(2):
            cx = left_x + Inches(0.6) + c * (cell_w + Inches(0.1))
            spec = cell_specs[r * 2 + c]
            tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, ry, cell_w, cell_h)
            tile.fill.solid(); tile.fill.fore_color.rgb = WHITE
            tile.line.color.rgb = spec[2]; tile.line.width = Pt(2)
            tile.shadow.inherit = False
            add_text(s, cx, ry + Inches(0.18), cell_w, Inches(0.35),
                     spec[0], size=11, bold=True, color=spec[2], align=PP_ALIGN.CENTER)
            add_text(s, cx, ry + Inches(0.55), cell_w, Inches(0.5),
                     spec[1], size=18, bold=True, color=INK, font=HEAD_FONT, align=PP_ALIGN.CENTER)

    # Right: Plain-English "what we feed in"
    right_x = Inches(7.0)
    add_text(s, right_x, top_y, Inches(6.0), Inches(0.4),
             "What each model looks at",
             size=15, bold=True, color=INK, font=HEAD_FONT)

    # In/Out box: simple list of inputs and one output
    in_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top_y + Inches(0.4),
                                Inches(6.0), Inches(2.6))
    in_box.fill.solid(); in_box.fill.fore_color.rgb = INK
    in_box.line.fill.background(); in_box.shadow.inherit = False

    add_text(s, right_x + Inches(0.25), top_y + Inches(0.5), Inches(5.6), Inches(0.3),
             "INPUTS (8 factors per month)", size=10, bold=True, color=AMBER, font=BODY_FONT)
    add_bullets(s, right_x + Inches(0.25), top_y + Inches(0.85), Inches(5.6), Inches(2.2), [
        "How much oil the US produced",
        "How much oil was in storage",
        "Last month's oil price",
        "Net exports of oil",
        "How busy refineries were",
        "Strength of the US dollar",
        "Global geopolitical risk score",
        "Whether the Strait of Hormuz was under threat",
    ], size=10, color=CREAM, bullet_color=AMBER, line_spacing=1.05)

    # Sources, simple language
    add_text(s, right_x, top_y + Inches(3.15), Inches(6.0), Inches(0.4),
             "Where the data comes from", size=15, bold=True, color=INK, font=HEAD_FONT)
    add_bullets(s, right_x, top_y + Inches(3.55), Inches(6.0), Inches(2.5), [
        ("Federal Reserve", "oil prices, inflation, US dollar index"),
        ("Energy Information Administration", "production, storage, refinery activity, exports"),
        ("Geopolitical Risk Index", "monthly score from a published academic source"),
        ("All prices adjusted for inflation", "shown in today's dollars (2025 base year)"),
    ], size=11, color=INK, line_spacing=1.15)

    add_slide_number(s, n, total)


def slide_forecast(prs, n, total):
    """Featured slide: 4-month forecast chart with detail."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, "4-month forecast through August 2026", kicker="The forecast")

    f = _read_forecast_summary()

    # Chart - left
    add_image(s, CHARTS / "07_forecast.png", Inches(0.55), Inches(1.5), w=Inches(8.4))

    # Right side: prediction summary tiles + interpretation
    rx = Inches(9.2)
    add_text(s, rx, Inches(1.5), Inches(4.0), Inches(0.4),
             "August 2026 outlook", size=15, bold=True, color=INK, font=HEAD_FONT)

    # Light tile
    light_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    rx, Inches(2.0), Inches(4.0), Inches(1.1))
    light_box.fill.solid(); light_box.fill.fore_color.rgb = WHITE
    light_box.line.color.rgb = TEAL; light_box.line.width = Pt(1.5)
    light_box.shadow.inherit = False
    add_text(s, rx + Inches(0.2), Inches(2.1), Inches(3.6), Inches(0.3),
             "LIGHT OIL (WTI)", size=10, bold=True, color=TEAL, font=BODY_FONT)
    add_text(s, rx + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.55),
             f"${f['light_end']:.0f}", size=32, bold=True, color=INK, font=HEAD_FONT)
    add_text(s, rx + Inches(1.55), Inches(2.55), Inches(2.3), Inches(0.4),
             f"from ${f['light_last']:.0f} today", size=11, color=GRAY, font=BODY_FONT)
    add_text(s, rx + Inches(0.2), Inches(2.85), Inches(3.6), Inches(0.25),
             f"95% range: ${f['light_lo']:.0f} \u2013 ${f['light_hi']:.0f}",
             size=10, italic=True, color=GRAY, font=BODY_FONT)

    # Heavy tile
    heavy_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    rx, Inches(3.2), Inches(4.0), Inches(1.1))
    heavy_box.fill.solid(); heavy_box.fill.fore_color.rgb = WHITE
    heavy_box.line.color.rgb = AMBER; heavy_box.line.width = Pt(1.5)
    heavy_box.shadow.inherit = False
    add_text(s, rx + Inches(0.2), Inches(3.3), Inches(3.6), Inches(0.3),
             "HEAVY OIL", size=10, bold=True, color=AMBER, font=BODY_FONT)
    add_text(s, rx + Inches(0.2), Inches(3.6), Inches(3.6), Inches(0.55),
             f"${f['heavy_end']:.0f}", size=32, bold=True, color=INK, font=HEAD_FONT)
    add_text(s, rx + Inches(1.55), Inches(3.75), Inches(2.3), Inches(0.4),
             f"from ${f['heavy_last']:.0f} today", size=11, color=GRAY, font=BODY_FONT)
    add_text(s, rx + Inches(0.2), Inches(4.05), Inches(3.6), Inches(0.25),
             f"95% range: ${f['heavy_lo']:.0f} \u2013 ${f['heavy_hi']:.0f}",
             size=10, italic=True, color=GRAY, font=BODY_FONT)

    # Interpretation bullets
    add_text(s, rx, Inches(4.45), Inches(4.0), Inches(0.4),
             "Why this story", size=13, bold=True, color=INK, font=HEAD_FONT)
    add_bullets(s, rx, Inches(4.85), Inches(4.0), Inches(2.3), [
        "Prices currently elevated by ongoing Israel-Iran conflict",
        "Model expects gradual mean reversion as the spike fades",
        "Heavy-light gap roughly stable around $5\u20136",
        "Confidence range widens with horizon \u2014 longer = less certain",
    ], size=11, color=INK, line_spacing=1.2, bullet_color=AMBER)

    # Method footnote at bottom
    add_text(s, Inches(0.55), Inches(7.0), Inches(8.5), Inches(0.3),
             "Method: long-conflict regression with supply/demand inputs held at latest values; "
             "seasonality recomputed each month.",
             size=9, italic=True, color=GRAY, font=BODY_FONT)

    add_slide_number(s, n, total)


def slide_chart(prs, n, total, *, kicker, title, chart_path: Path,
                interpretation: list, chart_w_in=8.4):
    """Generic content slide: chart left, interpretation right."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, title, kicker=kicker)
    # Chart
    add_image(s, chart_path, Inches(0.55), Inches(1.5), w=Inches(chart_w_in))
    # Interpretation bullets
    add_text(s, Inches(9.2), Inches(1.5), Inches(4.0), Inches(0.4),
             "Interpretation", size=15, bold=True, color=INK, font=HEAD_FONT)
    add_bullets(s, Inches(9.2), Inches(1.95), Inches(4.0), Inches(5.2),
                interpretation, size=12, color=INK, line_spacing=1.25)
    add_slide_number(s, n, total)


def slide_chart_full(prs, n, total, *, kicker, title, chart_path: Path,
                     caption: str, interpretation: list):
    """Chart on top, two-column interpretation below."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, title, kicker=kicker)
    add_image(s, chart_path, Inches(0.55), Inches(1.4), w=Inches(12.2))
    add_text(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.35),
             caption, size=11, italic=True, color=GRAY, font=BODY_FONT)

    half = len(interpretation) // 2 + len(interpretation) % 2
    add_bullets(s, Inches(0.55), Inches(6.25), Inches(6.0), Inches(1.1),
                interpretation[:half], size=11, color=INK, line_spacing=1.2)
    add_bullets(s, Inches(6.85), Inches(6.25), Inches(6.0), Inches(1.1),
                interpretation[half:], size=11, color=INK, line_spacing=1.2)
    add_slide_number(s, n, total)


def slide_limitations(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, "What to keep in mind", kicker="Honest limits")

    # Left column: data caveats
    add_text(s, Inches(0.55), Inches(1.5), Inches(6.2), Inches(0.4),
             "Data", size=18, bold=True, color=TEAL, font=HEAD_FONT)
    add_bullets(s, Inches(0.55), Inches(1.95), Inches(6.2), Inches(2.3), [
        ("Heavy oil price uses a stand-in",
            "we used a heavy oil blend that's similar to WCS, since clean WCS history is hard to source"),
        ("Reliable data only goes back to about 2006",
            "some sources started later; we use what's complete"),
        ("Hormuz threat status is a judgment call",
            "the strait has never actually closed, so we flagged periods of credible threat"),
        ("Risk index relies on an outside source",
            "if the source is offline, we fall back to a neutral value"),
    ], size=12, color=INK, line_spacing=1.2)

    # Right column: model + forecast caveats
    add_text(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.4),
             "Model & forecast", size=18, bold=True, color=RED, font=HEAD_FONT)
    add_bullets(s, Inches(7.0), Inches(1.95), Inches(6.0), Inches(3.2), [
        ("Forecast assumes today's conditions hold",
            "supply, demand, dollar, and conflict are held at latest values; a regime shift would change things"),
        ("Confidence range widens with horizon",
            "by August 2026 the 95% range spans roughly $30\u2013$110; near-term months are tighter"),
        ("Last month's price drives a lot of prediction",
            "realistic for oil, but means a sudden shock would push the whole forecast"),
        ("Some factors couldn't be measured cleanly",
            "Hormuz, geopolitics, refining margins \u2014 too rare or too small to detect"),
    ], size=12, color=INK, line_spacing=1.15)

    # Box at bottom for honest assessment
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.65),
                             Inches(12.4), Inches(1.3))
    box.fill.solid(); box.fill.fore_color.rgb = INK
    box.line.fill.background(); box.shadow.inherit = False
    add_text(s, Inches(0.85), Inches(5.75), Inches(11.8), Inches(0.4),
             "The honest take", size=13, bold=True, color=AMBER, font=BODY_FONT)
    add_text(s, Inches(0.85), Inches(6.1), Inches(11.8), Inches(0.85),
             "The forecast is a baseline, not a guarantee. It says: if today's geopolitics, supply, and dollar levels "
             "persist, prices drift back toward the model's long-conflict average. A new shock \u2014 escalation, a "
             "production cut, a recession \u2014 would shift the picture meaningfully.",
             size=12, italic=True, color=CREAM, font=BODY_FONT)
    add_slide_number(s, n, total)


def slide_recommendations(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    page_header(s, "What to do next", kicker="Path forward")

    items = [
        ("1. Tighten the confidence ranges",
            "A small statistical adjustment will give us more reliable error bars without changing the predictions themselves.",
            TEAL),
        ("2. Get real Western Canadian Select price history",
            "We're using a similar heavy-oil blend as a stand-in. Real WCS data would sharpen the heavy-oil models.",
            AMBER),
        ("3. Drop factors that don't measurably affect prices",
            "Hormuz, geopolitical risk, and refining margins didn't move the needle. Removing them simplifies the story.",
            RED),
        ("4. Test the model on data it hasn't seen",
            "Re-fit on 2003–2023 and check how well it predicts 2024–2026. This is the real test of forecasting power.",
            INK),
        ("5. Try a time-series model for long-conflict periods",
            "Because last month's price has so much weight, a model designed for time series may be a better fit.",
            GRAY),
    ]
    y = Inches(1.5)
    for title, body, color in items:
        # circle bullet
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), y + Inches(0.05),
                                  Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background(); circ.shadow.inherit = False
        # number
        add_text(s, Inches(0.55), y + Inches(0.05), Inches(0.55), Inches(0.55),
                 title.split(".")[0], size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
        # text block
        add_text(s, Inches(1.3), y, Inches(11.7), Inches(0.4),
                 title.split(".", 1)[1].strip(), size=14, bold=True, color=INK, font=HEAD_FONT)
        add_text(s, Inches(1.3), y + Inches(0.4), Inches(11.7), Inches(0.7),
                 body, size=11, color=GRAY, font=BODY_FONT)
        y += Inches(1.05)

    add_slide_number(s, n, total)


def slide_closing(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = AMBER; band.line.fill.background()

    add_text(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
             "Questions?",
             size=54, bold=True, color=WHITE, font=HEAD_FONT)
    add_text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.6),
             "Everything in this deck is reproducible from raw public data",
             size=18, italic=True, color=CREAM, font=BODY_FONT)

    add_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.4),
             "WHAT YOU CAN EXPLORE",
             size=11, bold=True, color=AMBER, font=BODY_FONT)
    add_bullets(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.6), [
        "A 4-month forecast table with prediction and confidence range per month",
        "A spreadsheet with all four model results and statistical details",
        "The seven charts shown in this deck, as image files",
        "The complete monthly dataset and the settings file used to build the forecast",
    ], size=13, color=CREAM, bullet_color=AMBER)


# -- Build ------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 12

    # 1
    slide_title(prs, total)
    # 2
    slide_executive_summary(prs, 2, total)
    # 3
    slide_methodology(prs, 3, total)
    # 4 - THE FORECAST (centerpiece)
    slide_forecast(prs, 4, total)
    # 5 - price history
    slide_chart(prs, 5, total,
        kicker="Chart 1 of 5",
        title="Oil prices through major Middle East conflicts",
        chart_path=CHARTS / "01_price_history_with_regimes.png",
        chart_w_in=8.4,
        interpretation=[
            ("2008 spike near $200 a barrel", "global demand was huge while supply was tight"),
            ("2014–2016 crash", "US shale boom flooded the market; the ISIS conflict didn't push prices up"),
            ("2020 COVID drop near $25", "demand collapsed when the world stopped moving"),
            ("Light and heavy oil track each other", "they tend to move together, just at different price levels"),
            ("Conflict periods don't always spike prices", "supply and demand often matter more than wars"),
            ("Big takeaway", "the economy moves oil more than geopolitics does"),
        ])
    # 6 - actual vs fitted
    slide_chart(prs, 6, total,
        kicker="Chart 2 of 5",
        title="How well the model matches reality",
        chart_path=CHARTS / "02_actual_vs_fitted.png",
        chart_w_in=8.4,
        interpretation=[
            ("Each dot is one month", "the closer to the diagonal line, the better the prediction"),
            ("All four models hug the line tightly", "the model captures the historic record well"),
            ("Heavy oil + long conflicts is the best fit", "96% of price moves explained"),
            ("Light oil + short conflicts has the most spread", "still a strong 92% fit"),
            ("Important caveat", "matching the past is easier than predicting the future"),
        ])
    # 7 - coefficients
    slide_chart(prs, 7, total,
        kicker="Chart 3 of 5",
        title="Which factors actually matter",
        chart_path=CHARTS / "05_coefficients_comparison.png",
        chart_w_in=8.4,
        interpretation=[
            ("Last month's price has the biggest effect", "oil prices have momentum; today builds on yesterday"),
            ("US dollar strength: strong negative impact", "when the dollar rises, oil priced in dollars tends to fall"),
            ("More US production = lower prices", "as expected; basic supply and demand"),
            ("Storage levels behave differently in long conflicts", "people build buffer stocks during sustained tensions"),
            ("Hormuz threats and risk index: surprisingly small", "their effect on prices is too small to measure reliably"),
            ("Bars show the range of uncertainty", "bars crossing zero mean the effect could go either way"),
        ])
    # 8 - diagnostics (heavy long-term, the best-fit cell)
    slide_chart(prs, 8, total,
        kicker="Chart 4 of 5",
        title="Quality check on the best model",
        chart_path=CHARTS / "04_diagnostics_heavy_longterm.png",
        chart_w_in=7.0,
        interpretation=[
            ("These four panels test whether the model is well-behaved", "they look at the leftover errors after each prediction"),
            ("Top-left and top-right: errors are roughly symmetric", "no systematic bias in either direction"),
            ("Bottom-left: errors look randomly scattered", "the model isn't missing an obvious pattern"),
            ("Bottom-right: a small lingering issue", "this month's error nudges next month's, slightly"),
            ("Practical impact", "predictions are accurate, but our confidence ranges are slightly too tight"),
            ("Easy fix exists", "switch to a more robust statistical method (one-line code change)"),
        ])
    # 9 - multicollinearity
    slide_chart(prs, 9, total,
        kicker="Chart 5 of 5",
        title="Which factors overlap with each other",
        chart_path=CHARTS / "06_correlation_heatmap.png",
        chart_w_in=8.4,
        interpretation=[
            ("This chart shows which factors move together", "deep red = move in lockstep, deep blue = move opposite"),
            ("Light and heavy oil prices: nearly identical movement", "that's why we split them into separate models"),
            ("Production and storage rose together (2010s shale era)", "they share some information but not enough to cause trouble"),
            ("Dollar strength is mostly independent", "good — its effect on oil is clean and easy to pinpoint"),
            ("Risk index doesn't track other factors", "explains why its impact on price is hard to nail down"),
            ("Big takeaway", "no harmful overlap — each factor brings its own information"),
        ])
    # 10 - limitations
    slide_limitations(prs, 10, total)
    # 11 - recommendations
    slide_recommendations(prs, 11, total)
    # 12 - closing
    slide_closing(prs, 12, total)

    prs.save(OUTFILE)
    print(f"Saved deck: {OUTFILE}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
